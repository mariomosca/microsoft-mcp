"""Unit tests for email attachment resolution (no network / no live account).

Covers the remote/HTTP deploy fix: email attachments must be passable inline as
base64 (``attachments_inline``) because the Azure container does not share a
filesystem with the client. Local-path mode (``attachments``) is kept for stdio.
"""

import base64

import pytest

from microsoft_mcp.tools import _resolve_email_attachments


def _b64(data: bytes) -> str:
    return base64.b64encode(data).decode("utf-8")


def test_none_returns_empty():
    assert _resolve_email_attachments(None, None) == []


def test_inline_single():
    payload = b"hello world"
    out = _resolve_email_attachments(
        None, [{"name": "doc.txt", "content_base64": _b64(payload)}]
    )
    assert out == [("doc.txt", payload)]


def test_inline_multiple_preserves_order_and_names():
    out = _resolve_email_attachments(
        None,
        [
            {"name": "a.pdf", "content_base64": _b64(b"AAA")},
            {"name": "b.pdf", "content_base64": _b64(b"BBBB")},
        ],
    )
    assert out == [("a.pdf", b"AAA"), ("b.pdf", b"BBBB")]


def test_local_path(tmp_path):
    f = tmp_path / "report.txt"
    f.write_bytes(b"local bytes")
    out = _resolve_email_attachments(str(f), None)
    assert out == [("report.txt", b"local bytes")]


def test_local_and_inline_combined(tmp_path):
    f = tmp_path / "local.bin"
    f.write_bytes(b"\x00\x01\x02")
    out = _resolve_email_attachments(
        str(f), [{"name": "inline.txt", "content_base64": _b64(b"xyz")}]
    )
    assert out == [("local.bin", b"\x00\x01\x02"), ("inline.txt", b"xyz")]


def test_inline_missing_name_raises():
    with pytest.raises(ValueError, match="name.*content_base64"):
        _resolve_email_attachments(None, [{"content_base64": _b64(b"x")}])


def test_inline_missing_content_raises():
    with pytest.raises(ValueError, match="name.*content_base64"):
        _resolve_email_attachments(None, [{"name": "x.txt"}])


def test_inline_invalid_base64_raises():
    with pytest.raises(ValueError, match="Invalid base64"):
        _resolve_email_attachments(
            None, [{"name": "x.txt", "content_base64": "not!valid!base64"}]
        )


# --- Calendar event attachments (mocked graph, no network) -----------------

from microsoft_mcp import tools  # noqa: E402


def test_list_event_attachments_maps_metadata(monkeypatch):
    captured = {}

    def fake_paginated(path, account_id=None, params=None, limit=None):
        captured["path"] = path
        captured["params"] = params
        yield {
            "id": "att-1",
            "name": "ticket.pdf",
            "contentType": "application/pdf",
            "size": 1234,
            "isInline": False,
        }

    monkeypatch.setattr(tools.graph, "request_paginated", fake_paginated)

    out = tools.list_event_attachments(event_id="evt-1", account_id="acc-1")

    assert captured["path"] == "/me/events/evt-1/attachments"
    assert out == [
        {
            "id": "att-1",
            "name": "ticket.pdf",
            "content_type": "application/pdf",
            "size": 1234,
            "is_inline": False,
        }
    ]


def _mock_attachment(monkeypatch, name, content_type, raw: bytes):
    b64 = base64.b64encode(raw).decode("utf-8")
    monkeypatch.setattr(
        tools.graph,
        "request",
        lambda *a, **k: {
            "name": name,
            "contentType": content_type,
            "size": len(raw),
            "contentBytes": b64,
        },
    )


# --- get_event_attachment: NEVER inline base64 over HTTP --------------------


def test_get_event_attachment_http_stages_to_onedrive(monkeypatch):
    raw = b"PNR-ABC123 flight 09:40"
    _mock_attachment(monkeypatch, "ticket.pdf", "application/pdf", raw)

    uploaded = {}

    def fake_upload(path, data, account_id=None, item_properties=None):
        uploaded["path"] = path
        uploaded["bytes"] = len(data)
        return {"id": "drive-item-1", "webUrl": "https://onedrive/ticket.pdf"}

    monkeypatch.setattr(tools.graph, "upload_large_file", fake_upload)

    out = tools.get_event_attachment(
        event_id="evt-1", attachment_id="att-1", account_id="acc-1"
    )

    # the key regression guard: bytes are NEVER returned inline
    assert "content_base64" not in out
    assert out["onedrive_file_id"] == "drive-item-1"
    assert out["web_url"] == "https://onedrive/ticket.pdf"
    assert out["onedrive_path"] == "Attachments/Events/ticket.pdf"
    assert uploaded["bytes"] == len(raw)
    assert uploaded["path"] == "/me/drive/root:/Attachments/Events/ticket.pdf:"


def test_get_event_attachment_save_path(monkeypatch, tmp_path):
    raw = b"binary-bytes"
    _mock_attachment(monkeypatch, "f.bin", "application/octet-stream", raw)

    dest = tmp_path / "out" / "f.bin"
    out = tools.get_event_attachment(
        event_id="evt-1",
        attachment_id="att-1",
        account_id="acc-1",
        save_path=str(dest),
    )

    assert dest.read_bytes() == raw
    assert out["saved_to"] == str(dest.resolve())
    assert "content_base64" not in out


def test_get_event_attachment_missing_content_raises(monkeypatch):
    monkeypatch.setattr(
        tools.graph,
        "request",
        lambda *a, **k: {"name": "f", "contentType": "x", "size": 0},
    )
    with pytest.raises(ValueError, match="reference attachment"):
        tools.get_event_attachment(
            event_id="e", attachment_id="a", account_id="acc"
        )


# --- read_event_attachment: server-side text extraction ---------------------


def test_read_event_attachment_xlsx(monkeypatch):
    import io
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Costi"
    ws.append(["Voce", "Importo"])
    ws.append(["Hotel", 120])
    ws.append(["Volo", 350])
    buf = io.BytesIO()
    wb.save(buf)
    raw = buf.getvalue()

    _mock_attachment(
        monkeypatch,
        "FORM SCHEDA COSTI.xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        raw,
    )

    out = tools.read_event_attachment(
        event_id="evt-1", attachment_id="att-1", account_id="acc-1"
    )
    assert out["kind"] == "xlsx"
    assert "content_base64" not in out
    assert "Costi" in out["text"]
    assert "Hotel" in out["text"]
    assert "350" in out["text"]


def test_read_event_attachment_csv(monkeypatch):
    raw = b"col1,col2\na,1\nb,2\n"
    _mock_attachment(monkeypatch, "data.csv", "text/csv", raw)

    out = tools.read_event_attachment(
        event_id="e", attachment_id="a", account_id="acc"
    )
    assert out["kind"] == "csv"
    assert out["text"] == "col1,col2\na,1\nb,2\n"


def test_read_event_attachment_text_truncates(monkeypatch):
    raw = ("x" * 100).encode()
    _mock_attachment(monkeypatch, "big.txt", "text/plain", raw)

    out = tools.read_event_attachment(
        event_id="e", attachment_id="a", account_id="acc", max_chars=10
    )
    assert out["truncated"] is True
    assert len(out["text"]) == 10


def test_read_event_attachment_unsupported_image(monkeypatch):
    raw = b"\x89PNG\r\n\x1a\n" + b"\x00" * 20
    _mock_attachment(monkeypatch, "photo.png", "image/png", raw)

    out = tools.read_event_attachment(
        event_id="e", attachment_id="a", account_id="acc"
    )
    assert out["kind"] == "unsupported"
    assert out["text"] is None
    assert "get_event_attachment" in out["note"]


# --- _extract_text: extended formats ----------------------------------------


def test_extract_html():
    raw = b"<html><body><h1>Titolo</h1><p>Ciao</p><script>x()</script></body></html>"
    kind, text = tools._extract_text("page.html", "text/html", raw)
    assert kind == "html"
    assert "Titolo" in text and "Ciao" in text
    assert "x()" not in text  # script stripped


def test_extract_rtf():
    raw = rb"{\rtf1\ansi Hello \b world\b0 .}"
    kind, text = tools._extract_text("note.rtf", "application/rtf", raw)
    assert kind == "rtf"
    assert "Hello" in text and "world" in text


def test_extract_pptx():
    import io
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Budget Q3"
    buf = io.BytesIO()
    prs.save(buf)

    kind, text = tools._extract_text("deck.pptx", "application/vnd.ms-powerpoint", buf.getvalue())
    assert kind == "pptx"
    assert "Budget Q3" in text


def test_extract_eml():
    raw = (
        b"From: a@x.com\r\nTo: b@y.com\r\nSubject: Test\r\n"
        b"Content-Type: text/plain\r\n\r\nCorpo della mail.\r\n"
    )
    kind, text = tools._extract_text("m.eml", "message/rfc822", raw)
    assert kind == "eml"
    assert "Subject: Test" in text
    assert "Corpo della mail." in text


def test_extract_unsupported_image():
    kind, text = tools._extract_text("x.png", "image/png", b"\x89PNG\r\n")
    assert kind == "unsupported"
    assert text == ""


# --- read_attachment_text: generic multi-source reader ----------------------


def test_read_attachment_text_requires_a_source():
    with pytest.raises(ValueError, match="Provide a source"):
        tools.read_attachment_text(account_id="acc")


def test_read_attachment_text_from_event(monkeypatch):
    raw = b"plain ticket PNR-XYZ"
    _mock_attachment(monkeypatch, "t.txt", "text/plain", raw)
    out = tools.read_attachment_text(
        account_id="acc", event_id="e", attachment_id="a"
    )
    assert out["kind"] == "text"
    assert "PNR-XYZ" in out["text"]


def test_read_attachment_text_from_onedrive(monkeypatch):
    raw = b"col1,col2\nx,1\n"

    monkeypatch.setattr(
        tools.graph,
        "request",
        lambda *a, **k: {
            "name": "data.csv",
            "file": {"mimeType": "text/csv"},
            "size": len(raw),
            "@microsoft.graph.downloadUrl": "https://dl/data.csv",
        },
    )

    class FakeResp:
        content = raw

        def raise_for_status(self):
            pass

    import httpx

    monkeypatch.setattr(httpx, "get", lambda *a, **k: FakeResp())

    out = tools.read_attachment_text(account_id="acc", onedrive_file_id="file-1")
    assert out["kind"] == "csv"
    assert out["text"] == "col1,col2\nx,1\n"
