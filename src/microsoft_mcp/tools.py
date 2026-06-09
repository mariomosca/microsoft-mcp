import base64
import datetime as dt
import pathlib as pl
from typing import Any
from fastmcp import FastMCP
from . import graph, auth

mcp = FastMCP("microsoft-mcp")

FOLDERS = {
    k.casefold(): v
    for k, v in {
        "inbox": "inbox",
        "sent": "sentitems",
        "drafts": "drafts",
        "deleted": "deleteditems",
        "junk": "junkemail",
        "archive": "archive",
    }.items()
}


@mcp.tool
def list_accounts() -> list[dict[str, str]]:
    """List all signed-in Microsoft accounts"""
    from .auth_context import current_graph_token

    # HTTP mode: user is already authenticated via OAuth/OBO
    if current_graph_token.get() is not None:
        try:
            from fastmcp.server.dependencies import get_access_token

            token = get_access_token()
            if token and hasattr(token, "claims"):
                upstream = (token.claims or {}).get("upstream_claims", {})
                return [
                    {
                        "username": upstream.get("upn", upstream.get("name", "authenticated-user")),
                        "account_id": upstream.get("oid", "http-user"),
                    }
                ]
        except Exception:
            pass

    # Stdio mode: return from MSAL cache
    return [
        {"username": acc.username, "account_id": acc.account_id}
        for acc in auth.list_accounts()
    ]


@mcp.tool
def authenticate_account() -> dict[str, str]:
    """Authenticate a new Microsoft account using device flow authentication

    Returns authentication instructions and device code for the user to complete authentication.
    The user must visit the URL and enter the code to authenticate their Microsoft account.
    """
    app = auth.get_app()
    flow = app.initiate_device_flow(scopes=auth.SCOPES)

    if "user_code" not in flow:
        error_msg = flow.get("error_description", "Unknown error")
        raise Exception(f"Failed to get device code: {error_msg}")

    verification_url = flow.get(
        "verification_uri",
        flow.get("verification_url", "https://microsoft.com/devicelogin"),
    )

    return {
        "status": "authentication_required",
        "instructions": "To authenticate a new Microsoft account:",
        "step1": f"Visit: {verification_url}",
        "step2": f"Enter code: {flow['user_code']}",
        "step3": "Sign in with the Microsoft account you want to add",
        "step4": "After authenticating, use the 'complete_authentication' tool to finish the process",
        "device_code": flow["user_code"],
        "verification_url": verification_url,
        "expires_in": str(flow.get("expires_in", 900)),
        "_flow_cache": str(flow),
    }


@mcp.tool
def complete_authentication(flow_cache: str) -> dict[str, str]:
    """Complete the authentication process after the user has entered the device code

    Args:
        flow_cache: The flow data returned from authenticate_account (the _flow_cache field)

    Returns:
        Account information if authentication was successful
    """
    import ast

    try:
        flow = ast.literal_eval(flow_cache)
    except (ValueError, SyntaxError):
        raise ValueError("Invalid flow cache data")

    app = auth.get_app()
    result = app.acquire_token_by_device_flow(flow)

    if "error" in result:
        error_msg = result.get("error_description", result["error"])
        if "authorization_pending" in error_msg:
            return {
                "status": "pending",
                "message": "Authentication is still pending. The user needs to complete the authentication process.",
                "instructions": "Please ensure you've visited the URL and entered the code, then try again.",
            }
        raise Exception(f"Authentication failed: {error_msg}")

    # Save the token cache
    cache = app.token_cache
    if isinstance(cache, auth.msal.SerializableTokenCache) and cache.has_state_changed:
        auth._write_cache(cache.serialize())

    # Get the newly added account
    accounts = app.get_accounts()
    if accounts:
        # Find the account that matches the token we just got
        for account in accounts:
            if (
                account.get("username", "").lower()
                == result.get("id_token_claims", {})
                .get("preferred_username", "")
                .lower()
            ):
                return {
                    "status": "success",
                    "username": account["username"],
                    "account_id": account["home_account_id"],
                    "message": f"Successfully authenticated {account['username']}",
                }
        # If exact match not found, return the last account
        account = accounts[-1]
        return {
            "status": "success",
            "username": account["username"],
            "account_id": account["home_account_id"],
            "message": f"Successfully authenticated {account['username']}",
        }

    return {
        "status": "error",
        "message": "Authentication succeeded but no account was found",
    }


@mcp.tool
def list_emails(
    account_id: str,
    folder: str = "inbox",
    limit: int = 10,
    include_body: bool = True,
) -> list[dict[str, Any]]:
    """List emails from specified folder"""
    folder_path = FOLDERS.get(folder.casefold(), folder)

    if include_body:
        select_fields = "id,subject,from,toRecipients,ccRecipients,receivedDateTime,hasAttachments,body,conversationId,isRead"
    else:
        select_fields = "id,subject,from,toRecipients,receivedDateTime,hasAttachments,conversationId,isRead"

    params = {
        "$top": min(limit, 100),
        "$select": select_fields,
        "$orderby": "receivedDateTime desc",
    }

    emails = list(
        graph.request_paginated(
            f"/me/mailFolders/{folder_path}/messages",
            account_id,
            params=params,
            limit=limit,
        )
    )

    return emails


@mcp.tool
def get_email(
    email_id: str,
    account_id: str,
    include_body: bool = True,
    body_max_length: int = 50000,
    include_attachments: bool = True,
) -> dict[str, Any]:
    """Get email details with size limits

    Args:
        email_id: The email ID
        account_id: The account ID
        include_body: Whether to include the email body (default: True)
        body_max_length: Maximum characters for body content (default: 50000)
        include_attachments: Whether to include attachment metadata (default: True)
    """
    params = {}
    if include_attachments:
        params["$expand"] = "attachments($select=id,name,size,contentType)"

    result = graph.request("GET", f"/me/messages/{email_id}", account_id, params=params)
    if not result:
        raise ValueError(f"Email with ID {email_id} not found")

    # Truncate body if needed
    if include_body and "body" in result and "content" in result["body"]:
        content = result["body"]["content"]
        if len(content) > body_max_length:
            result["body"]["content"] = (
                content[:body_max_length]
                + f"\n\n[Content truncated - {len(content)} total characters]"
            )
            result["body"]["truncated"] = True
            result["body"]["total_length"] = len(content)
    elif not include_body and "body" in result:
        del result["body"]

    # Remove attachment content bytes to reduce size
    if "attachments" in result and result["attachments"]:
        for attachment in result["attachments"]:
            if "contentBytes" in attachment:
                del attachment["contentBytes"]

    return result


@mcp.tool
def create_email_draft(
    account_id: str,
    to: str | list[str],
    subject: str,
    body: str,
    cc: str | list[str] | None = None,
    attachments: str | list[str] | None = None,
    attachments_inline: list[dict[str, str]] | None = None,
) -> dict[str, Any]:
    """Create an email draft, optionally with attachments.

    Attachments may be provided as local file path(s) (``attachments``, server
    filesystem only — works for local stdio) and/or inline as base64
    (``attachments_inline``, a list of ``{"name": str, "content_base64": str}``,
    required for remote/HTTP deployments where the server cannot read the
    client's filesystem). Both may be combined.
    """
    to_list = [to] if isinstance(to, str) else to

    message = {
        "subject": subject,
        "body": {"contentType": "Text", "content": body},
        "toRecipients": [{"emailAddress": {"address": addr}} for addr in to_list],
    }

    if cc:
        cc_list = [cc] if isinstance(cc, str) else cc
        message["ccRecipients"] = [
            {"emailAddress": {"address": addr}} for addr in cc_list
        ]

    small_attachments = []
    large_attachments = []

    for att_name, content_bytes in _resolve_email_attachments(
        attachments, attachments_inline
    ):
        att_size = len(content_bytes)
        if att_size < 3 * 1024 * 1024:
            small_attachments.append(
                {
                    "@odata.type": "#microsoft.graph.fileAttachment",
                    "name": att_name,
                    "contentBytes": base64.b64encode(content_bytes).decode("utf-8"),
                }
            )
        else:
            large_attachments.append(
                {
                    "name": att_name,
                    "content_bytes": content_bytes,
                    "content_type": "application/octet-stream",
                }
            )

    if small_attachments:
        message["attachments"] = small_attachments

    result = graph.request("POST", "/me/messages", account_id, json=message)
    if not result:
        raise ValueError("Failed to create email draft")

    message_id = result["id"]

    for att in large_attachments:
        graph.upload_large_mail_attachment(
            message_id,
            att["name"],
            att["content_bytes"],
            account_id,
            att.get("content_type", "application/octet-stream"),
        )

    return result


@mcp.tool
def send_email(
    account_id: str,
    to: str | list[str],
    subject: str,
    body: str,
    cc: str | list[str] | None = None,
    attachments: str | list[str] | None = None,
    attachments_inline: list[dict[str, str]] | None = None,
) -> dict[str, str]:
    """Send an email immediately, optionally with attachments.

    Attachments may be provided as local file path(s) (``attachments``, server
    filesystem only — works for local stdio) and/or inline as base64
    (``attachments_inline``, a list of ``{"name": str, "content_base64": str}``,
    required for remote/HTTP deployments where the server cannot read the
    client's filesystem). Both may be combined.
    """
    to_list = [to] if isinstance(to, str) else to

    message = {
        "subject": subject,
        "body": {"contentType": "Text", "content": body},
        "toRecipients": [{"emailAddress": {"address": addr}} for addr in to_list],
    }

    if cc:
        cc_list = [cc] if isinstance(cc, str) else cc
        message["ccRecipients"] = [
            {"emailAddress": {"address": addr}} for addr in cc_list
        ]

    # Check if we have large attachments
    has_large_attachments = False
    processed_attachments = []

    for att_name, content_bytes in _resolve_email_attachments(
        attachments, attachments_inline
    ):
        att_size = len(content_bytes)
        processed_attachments.append(
            {
                "name": att_name,
                "content_bytes": content_bytes,
                "content_type": "application/octet-stream",
                "size": att_size,
            }
        )
        if att_size >= 3 * 1024 * 1024:
            has_large_attachments = True

    if not has_large_attachments and processed_attachments:
        message["attachments"] = [
            {
                "@odata.type": "#microsoft.graph.fileAttachment",
                "name": att["name"],
                "contentBytes": base64.b64encode(att["content_bytes"]).decode("utf-8"),
            }
            for att in processed_attachments
        ]
        graph.request("POST", "/me/sendMail", account_id, json={"message": message})
        return {"status": "sent"}
    elif has_large_attachments:
        # Create draft first, then add large attachments, then send
        # We need to handle large attachments manually here
        to_list = [to] if isinstance(to, str) else to
        message = {
            "subject": subject,
            "body": {"contentType": "Text", "content": body},
            "toRecipients": [{"emailAddress": {"address": addr}} for addr in to_list],
        }
        if cc:
            cc_list = [cc] if isinstance(cc, str) else cc
            message["ccRecipients"] = [
                {"emailAddress": {"address": addr}} for addr in cc_list
            ]

        result = graph.request("POST", "/me/messages", account_id, json=message)
        if not result:
            raise ValueError("Failed to create email draft")

        message_id = result["id"]

        for att in processed_attachments:
            if att["size"] >= 3 * 1024 * 1024:
                graph.upload_large_mail_attachment(
                    message_id,
                    att["name"],
                    att["content_bytes"],
                    account_id,
                    att.get("content_type", "application/octet-stream"),
                )
            else:
                small_att = {
                    "@odata.type": "#microsoft.graph.fileAttachment",
                    "name": att["name"],
                    "contentBytes": base64.b64encode(att["content_bytes"]).decode(
                        "utf-8"
                    ),
                }
                graph.request(
                    "POST",
                    f"/me/messages/{message_id}/attachments",
                    account_id,
                    json=small_att,
                )

        graph.request("POST", f"/me/messages/{message_id}/send", account_id)
        return {"status": "sent"}
    else:
        graph.request("POST", "/me/sendMail", account_id, json={"message": message})
        return {"status": "sent"}


@mcp.tool
def update_email(
    email_id: str, updates: dict[str, Any], account_id: str
) -> dict[str, Any]:
    """Update email properties (isRead, categories, flag, etc.)"""
    result = graph.request(
        "PATCH", f"/me/messages/{email_id}", account_id, json=updates
    )
    if not result:
        raise ValueError(f"Failed to update email {email_id} - no response")
    return result


@mcp.tool
def delete_email(email_id: str, account_id: str) -> dict[str, str]:
    """Delete an email"""
    graph.request("DELETE", f"/me/messages/{email_id}", account_id)
    return {"status": "deleted"}


@mcp.tool
def move_email(
    email_id: str, destination_folder: str, account_id: str
) -> dict[str, Any]:
    """Move email to another folder"""
    folder_path = FOLDERS.get(destination_folder.casefold(), destination_folder)

    folders = graph.request("GET", "/me/mailFolders", account_id)
    folder_id = None

    if not folders:
        raise ValueError("Failed to retrieve mail folders")
    if "value" not in folders:
        raise ValueError(f"Unexpected folder response structure: {folders}")

    for folder in folders["value"]:
        if folder["displayName"].lower() == folder_path.lower():
            folder_id = folder["id"]
            break

    if not folder_id:
        raise ValueError(f"Folder '{destination_folder}' not found")

    payload = {"destinationId": folder_id}
    result = graph.request(
        "POST", f"/me/messages/{email_id}/move", account_id, json=payload
    )
    if not result:
        raise ValueError("Failed to move email - no response from server")
    if "id" not in result:
        raise ValueError(f"Failed to move email - unexpected response: {result}")
    return {"status": "moved", "new_id": result["id"]}


@mcp.tool
def reply_to_email(account_id: str, email_id: str, body: str) -> dict[str, str]:
    """Reply to an email (sender only)"""
    endpoint = f"/me/messages/{email_id}/reply"
    payload = {"message": {"body": {"contentType": "Text", "content": body}}}
    graph.request("POST", endpoint, account_id, json=payload)
    return {"status": "sent"}


@mcp.tool
def reply_all_email(account_id: str, email_id: str, body: str) -> dict[str, str]:
    """Reply to all recipients of an email"""
    endpoint = f"/me/messages/{email_id}/replyAll"
    payload = {"message": {"body": {"contentType": "Text", "content": body}}}
    graph.request("POST", endpoint, account_id, json=payload)
    return {"status": "sent"}


@mcp.tool
def create_reply_draft(
    account_id: str, email_id: str, body: str | None = None
) -> dict[str, Any]:
    """Create a reply draft (sender only) without sending. Returns the draft message for review."""
    endpoint = f"/me/messages/{email_id}/createReply"
    payload = {}
    if body:
        payload["message"] = {"body": {"contentType": "Text", "content": body}}
    result = graph.request("POST", endpoint, account_id, json=payload)
    if not result:
        raise ValueError("Failed to create reply draft")
    return result


@mcp.tool
def create_reply_all_draft(
    account_id: str, email_id: str, body: str | None = None
) -> dict[str, Any]:
    """Create a reply-all draft without sending. Preserves the email chain and all recipients. Returns the draft message for review."""
    endpoint = f"/me/messages/{email_id}/createReplyAll"
    payload = {}
    if body:
        payload["message"] = {"body": {"contentType": "Text", "content": body}}
    result = graph.request("POST", endpoint, account_id, json=payload)
    if not result:
        raise ValueError("Failed to create reply-all draft")
    return result


@mcp.tool
def list_events(
    account_id: str,
    days_ahead: int = 7,
    days_back: int = 0,
    include_details: bool = True,
) -> list[dict[str, Any]]:
    """List calendar events within specified date range, including recurring event instances"""
    now = dt.datetime.now(dt.timezone.utc)
    today_start = now.replace(hour=0, minute=0, second=0, microsecond=0)
    start = (today_start - dt.timedelta(days=days_back)).isoformat()
    end = (today_start + dt.timedelta(days=days_ahead + 1)).isoformat()

    params = {
        "startDateTime": start,
        "endDateTime": end,
        "$orderby": "start/dateTime",
        "$top": 100,
    }

    if include_details:
        params["$select"] = (
            "id,subject,start,end,location,body,attendees,organizer,isAllDay,recurrence,onlineMeeting,seriesMasterId"
        )
    else:
        params["$select"] = "id,subject,start,end,location,organizer,seriesMasterId"

    # Use calendarView to get recurring event instances
    events = list(
        graph.request_paginated("/me/calendarView", account_id, params=params)
    )

    return events


@mcp.tool
def get_event(event_id: str, account_id: str) -> dict[str, Any]:
    """Get full event details"""
    result = graph.request("GET", f"/me/events/{event_id}", account_id)
    if not result:
        raise ValueError(f"Event with ID {event_id} not found")
    return result


@mcp.tool
def create_event(
    account_id: str,
    subject: str,
    start: str,
    end: str,
    location: str | None = None,
    body: str | None = None,
    attendees: str | list[str] | None = None,
    timezone: str = "UTC",
    is_online_meeting: bool = False,
    online_meeting_provider: str = "teamsForBusiness",
    importance: str = "normal",
    is_all_day: bool = False,
    categories: list[str] | None = None,
    sensitivity: str = "normal",
    show_as: str = "busy",
    reminder_minutes: int = 15,
) -> dict[str, Any]:
    """Create a calendar event.

    Args:
        account_id: Microsoft account ID
        subject: Event title
        start: Start datetime (ISO format, e.g. 2026-04-10T14:30:00)
        end: End datetime (ISO format)
        location: Location display name
        body: Event body/description text
        attendees: Email address(es) of attendees
        timezone: Timezone (default UTC, e.g. Europe/Rome)
        is_online_meeting: If true, creates a Teams meeting link
        online_meeting_provider: Provider for online meeting (default teamsForBusiness)
        importance: Event importance (low, normal, high)
        is_all_day: If true, creates an all-day event
        categories: List of category names
        sensitivity: Event sensitivity (normal, personal, private, confidential)
        show_as: Free/busy status (free, tentative, busy, oof, workingElsewhere, unknown)
        reminder_minutes: Minutes before event to show reminder (default 15)
    """
    event = {
        "subject": subject,
        "start": {"dateTime": start, "timeZone": timezone},
        "end": {"dateTime": end, "timeZone": timezone},
        "importance": importance,
        "isAllDay": is_all_day,
        "sensitivity": sensitivity,
        "showAs": show_as,
        "isReminderOn": reminder_minutes > 0,
        "reminderMinutesBeforeStart": reminder_minutes,
    }

    if is_online_meeting:
        event["isOnlineMeeting"] = True
        event["onlineMeetingProvider"] = online_meeting_provider

    if location:
        event["location"] = {"displayName": location}

    if body:
        event["body"] = {"contentType": "Text", "content": body}

    if attendees:
        attendees_list = [attendees] if isinstance(attendees, str) else attendees
        event["attendees"] = [
            {"emailAddress": {"address": a}, "type": "required"} for a in attendees_list
        ]

    if categories:
        event["categories"] = categories

    result = graph.request("POST", "/me/events", account_id, json=event)
    if not result:
        raise ValueError("Failed to create event")
    return result


@mcp.tool
def update_event(
    event_id: str, updates: dict[str, Any], account_id: str
) -> dict[str, Any]:
    """Update event properties"""
    formatted_updates = {}

    if "subject" in updates:
        formatted_updates["subject"] = updates["subject"]
    if "start" in updates:
        formatted_updates["start"] = {
            "dateTime": updates["start"],
            "timeZone": updates.get("timezone", "UTC"),
        }
    if "end" in updates:
        formatted_updates["end"] = {
            "dateTime": updates["end"],
            "timeZone": updates.get("timezone", "UTC"),
        }
    if "location" in updates:
        formatted_updates["location"] = {"displayName": updates["location"]}
    if "body" in updates:
        formatted_updates["body"] = {"contentType": "Text", "content": updates["body"]}
    if "categories" in updates:
        formatted_updates["categories"] = updates["categories"]

    # Pass through any other Graph API properties not explicitly handled
    known_keys = {"subject", "start", "end", "location", "body", "categories", "timezone"}
    for key in updates:
        if key not in known_keys and key not in formatted_updates:
            formatted_updates[key] = updates[key]

    result = graph.request(
        "PATCH", f"/me/events/{event_id}", account_id, json=formatted_updates
    )
    return result or {"status": "updated"}


@mcp.tool
def delete_event(
    account_id: str, event_id: str, send_cancellation: bool = True
) -> dict[str, str]:
    """Delete or cancel a calendar event"""
    if send_cancellation:
        graph.request("POST", f"/me/events/{event_id}/cancel", account_id, json={})
    else:
        graph.request("DELETE", f"/me/events/{event_id}", account_id)
    return {"status": "deleted"}


@mcp.tool
def respond_event(
    account_id: str,
    event_id: str,
    response: str = "accept",
    message: str | None = None,
) -> dict[str, str]:
    """Respond to event invitation (accept, decline, tentativelyAccept)"""
    payload: dict[str, Any] = {"sendResponse": True}
    if message:
        payload["comment"] = message

    graph.request("POST", f"/me/events/{event_id}/{response}", account_id, json=payload)
    return {"status": response}


@mcp.tool
def check_availability(
    account_id: str,
    start: str,
    end: str,
    attendees: str | list[str] | None = None,
) -> dict[str, Any]:
    """Check calendar availability for scheduling"""
    me_info = graph.request("GET", "/me", account_id)
    if not me_info or "mail" not in me_info:
        raise ValueError("Failed to get user email address")
    schedules = [me_info["mail"]]
    if attendees:
        attendees_list = [attendees] if isinstance(attendees, str) else attendees
        schedules.extend(attendees_list)

    payload = {
        "schedules": schedules,
        "startTime": {"dateTime": start, "timeZone": "UTC"},
        "endTime": {"dateTime": end, "timeZone": "UTC"},
        "availabilityViewInterval": 30,
    }

    result = graph.request("POST", "/me/calendar/getSchedule", account_id, json=payload)
    if not result:
        raise ValueError("Failed to check availability")
    return result


@mcp.tool
def list_contacts(account_id: str, limit: int = 50) -> list[dict[str, Any]]:
    """List contacts"""
    params = {"$top": min(limit, 100)}

    contacts = list(
        graph.request_paginated("/me/contacts", account_id, params=params, limit=limit)
    )

    return contacts


@mcp.tool
def get_contact(contact_id: str, account_id: str) -> dict[str, Any]:
    """Get contact details"""
    result = graph.request("GET", f"/me/contacts/{contact_id}", account_id)
    if not result:
        raise ValueError(f"Contact with ID {contact_id} not found")
    return result


@mcp.tool
def create_contact(
    account_id: str,
    given_name: str,
    surname: str | None = None,
    email_addresses: str | list[str] | None = None,
    phone_numbers: dict[str, str] | None = None,
) -> dict[str, Any]:
    """Create a new contact"""
    contact: dict[str, Any] = {"givenName": given_name}

    if surname:
        contact["surname"] = surname

    if email_addresses:
        email_list = (
            [email_addresses] if isinstance(email_addresses, str) else email_addresses
        )
        contact["emailAddresses"] = [
            {"address": email, "name": f"{given_name} {surname or ''}".strip()}
            for email in email_list
        ]

    if phone_numbers:
        if "business" in phone_numbers:
            contact["businessPhones"] = [phone_numbers["business"]]
        if "home" in phone_numbers:
            contact["homePhones"] = [phone_numbers["home"]]
        if "mobile" in phone_numbers:
            contact["mobilePhone"] = phone_numbers["mobile"]

    result = graph.request("POST", "/me/contacts", account_id, json=contact)
    if not result:
        raise ValueError("Failed to create contact")
    return result


@mcp.tool
def update_contact(
    contact_id: str, updates: dict[str, Any], account_id: str
) -> dict[str, Any]:
    """Update contact information"""
    result = graph.request(
        "PATCH", f"/me/contacts/{contact_id}", account_id, json=updates
    )
    return result or {"status": "updated"}


@mcp.tool
def delete_contact(contact_id: str, account_id: str) -> dict[str, str]:
    """Delete a contact"""
    graph.request("DELETE", f"/me/contacts/{contact_id}", account_id)
    return {"status": "deleted"}


@mcp.tool
def list_files(
    account_id: str, path: str = "/", limit: int = 50
) -> list[dict[str, Any]]:
    """List files and folders in OneDrive"""
    endpoint = (
        "/me/drive/root/children"
        if path == "/"
        else f"/me/drive/root:/{path}:/children"
    )
    params = {
        "$top": min(limit, 100),
        "$select": "id,name,size,lastModifiedDateTime,folder,file,@microsoft.graph.downloadUrl",
    }

    items = list(
        graph.request_paginated(endpoint, account_id, params=params, limit=limit)
    )

    return [
        {
            "id": item["id"],
            "name": item["name"],
            "type": "folder" if "folder" in item else "file",
            "size": item.get("size", 0),
            "modified": item.get("lastModifiedDateTime"),
            "download_url": item.get("@microsoft.graph.downloadUrl"),
        }
        for item in items
    ]


@mcp.tool
def get_file(file_id: str, account_id: str, download_path: str) -> dict[str, Any]:
    """Download a file from OneDrive to local path"""
    import subprocess

    metadata = graph.request("GET", f"/me/drive/items/{file_id}", account_id)
    if not metadata:
        raise ValueError(f"File with ID {file_id} not found")

    download_url = metadata.get("@microsoft.graph.downloadUrl")
    if not download_url:
        raise ValueError("No download URL available for this file")

    try:
        subprocess.run(
            ["curl", "-L", "-o", download_path, download_url],
            check=True,
            capture_output=True,
        )

        return {
            "path": download_path,
            "name": metadata.get("name", "unknown"),
            "size_mb": round(metadata.get("size", 0) / (1024 * 1024), 2),
            "mime_type": metadata.get("file", {}).get("mimeType") if metadata else None,
        }
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"Failed to download file: {e.stderr.decode()}")


def _resolve_file_bytes(
    content_base64: str | None, local_file_path: str | None
) -> bytes:
    """Resolve file content from inline base64 or a local path.

    ``content_base64`` takes precedence. Raises if neither is provided.
    """
    if content_base64 is not None:
        try:
            return base64.b64decode(content_base64, validate=True)
        except (ValueError, base64.binascii.Error) as e:
            raise ValueError(f"Invalid base64 content: {e}")
    if local_file_path is not None:
        return pl.Path(local_file_path).expanduser().resolve().read_bytes()
    raise ValueError("Provide either content_base64 or local_file_path")


def _resolve_email_attachments(
    attachments: str | list[str] | None,
    attachments_inline: list[dict[str, str]] | None,
) -> list[tuple[str, bytes]]:
    """Resolve email attachments from local paths and/or inline base64.

    ``attachments`` are local file paths (server-side filesystem, stdio only).
    ``attachments_inline`` is a list of ``{"name": str, "content_base64": str}``
    objects, required for remote/HTTP deployments where the server cannot read
    the client's filesystem. Both may be combined. Returns ``(name, bytes)``.
    """
    resolved: list[tuple[str, bytes]] = []

    if attachments:
        paths = [attachments] if isinstance(attachments, str) else attachments
        for file_path in paths:
            path = pl.Path(file_path).expanduser().resolve()
            resolved.append((path.name, path.read_bytes()))

    if attachments_inline:
        for att in attachments_inline:
            name = att.get("name")
            content_base64 = att.get("content_base64")
            if not name or content_base64 is None:
                raise ValueError(
                    "Each inline attachment needs 'name' and 'content_base64'"
                )
            resolved.append((name, _resolve_file_bytes(content_base64, None)))

    return resolved


@mcp.tool
def create_file(
    onedrive_path: str,
    account_id: str,
    content_base64: str | None = None,
    local_file_path: str | None = None,
) -> dict[str, Any]:
    """Upload a file to OneDrive.

    Provide the content either inline as base64 (``content_base64``, required
    for remote/HTTP deployments where the server cannot read the client's
    filesystem) or as a path to a local file (``local_file_path``, only works
    when the server shares a filesystem with the caller, e.g. local stdio).
    """
    data = _resolve_file_bytes(content_base64, local_file_path)
    result = graph.upload_large_file(
        f"/me/drive/root:/{onedrive_path}:", data, account_id
    )
    if not result:
        raise ValueError(f"Failed to create file at path: {onedrive_path}")
    return result


@mcp.tool
def update_file(
    file_id: str,
    account_id: str,
    content_base64: str | None = None,
    local_file_path: str | None = None,
) -> dict[str, Any]:
    """Update OneDrive file content.

    Provide the new content either inline as base64 (``content_base64``,
    required for remote/HTTP deployments) or as a path to a local file
    (``local_file_path``, local stdio only).
    """
    data = _resolve_file_bytes(content_base64, local_file_path)
    result = graph.upload_large_file(f"/me/drive/items/{file_id}", data, account_id)
    if not result:
        raise ValueError(f"Failed to update file with ID: {file_id}")
    return result


@mcp.tool
def delete_file(file_id: str, account_id: str) -> dict[str, str]:
    """Delete a file or folder"""
    graph.request("DELETE", f"/me/drive/items/{file_id}", account_id)
    return {"status": "deleted"}


@mcp.tool
def get_attachment(
    email_id: str,
    attachment_id: str,
    account_id: str,
    save_path: str | None = None,
) -> dict[str, Any]:
    """Download an email attachment.

    If ``save_path`` is provided, the attachment is written to that path on the
    server's filesystem (only useful for local stdio). Otherwise the content is
    returned inline as base64 in ``content_base64`` - required for remote/HTTP
    deployments and ready to be passed back to ``create_file(content_base64=...)``.
    """
    result = graph.request(
        "GET", f"/me/messages/{email_id}/attachments/{attachment_id}", account_id
    )

    if not result:
        raise ValueError("Attachment not found")

    if "contentBytes" not in result:
        raise ValueError("Attachment content not available")

    response: dict[str, Any] = {
        "name": result.get("name", "unknown"),
        "content_type": result.get("contentType", "application/octet-stream"),
        "size": result.get("size", 0),
    }

    if save_path is not None:
        # Local stdio: persist to the server filesystem.
        path = pl.Path(save_path).expanduser().resolve()
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(base64.b64decode(result["contentBytes"]))
        response["saved_to"] = str(path)
    else:
        # Remote/HTTP: return the bytes inline (already base64 from Graph).
        response["content_base64"] = result["contentBytes"]

    return response


@mcp.tool
def list_event_attachments(
    event_id: str,
    account_id: str,
) -> list[dict[str, Any]]:
    """List attachments of a calendar event.

    Mirrors the email attachment pattern: returns metadata only (id, name,
    contentType, size) for each attachment so the content can be fetched on
    demand with ``get_event_attachment``. Useful e.g. to find a flight ticket
    (PDF/PNG) attached to a meeting and then extract PNR / times from it.
    """
    items = list(
        graph.request_paginated(
            f"/me/events/{event_id}/attachments",
            account_id,
            params={"$select": "id,name,contentType,size,isInline"},
        )
    )

    return [
        {
            "id": item.get("id"),
            "name": item.get("name", "unknown"),
            "content_type": item.get("contentType", "application/octet-stream"),
            "size": item.get("size", 0),
            "is_inline": item.get("isInline", False),
        }
        for item in items
    ]


def _fetch_event_attachment_bytes(
    event_id: str, attachment_id: str, account_id: str
) -> tuple[str, str, bytes]:
    """Fetch a single event fileAttachment. Returns (name, content_type, raw bytes).

    Raises ValueError for missing or reference (link) attachments.
    """
    result = graph.request(
        "GET",
        f"/me/events/{event_id}/attachments/{attachment_id}",
        account_id,
    )
    if not result:
        raise ValueError("Attachment not found")
    if "contentBytes" not in result:
        raise ValueError(
            "Attachment content not available - this is likely a reference "
            "attachment (link), which has no downloadable bytes. Open it from "
            "the link instead."
        )
    return (
        result.get("name", "unknown"),
        result.get("contentType", "application/octet-stream"),
        base64.b64decode(result["contentBytes"]),
    )


# How much extracted text to return at most, to avoid flooding the context.
ATTACHMENT_TEXT_MAX_CHARS = 50_000


def _extract_text(name: str, content_type: str, raw: bytes) -> tuple[str, str]:
    """Best-effort server-side text extraction. Returns (kind, text).

    Never returns binary to the model: the file is parsed to plain text here so
    the client can read it without a base64 round-trip. ``kind`` is a short
    label (text/csv/xlsx/pdf/docx/pptx/rtf/odf/html/eml/msg/unsupported).

    Supported today:
      text/csv/tsv/md/json/xml/log/yaml + any text/*  -> decoded
      xlsx/xlsm/ods                                   -> cells as TSV
      pdf                                             -> text per page
      docx/odt                                        -> paragraph text
      pptx                                            -> slide text
      rtf                                             -> stripped text
      html/htm                                        -> visible text
      eml                                             -> headers + body
      msg (Outlook)                                   -> headers + body

    Unsupported (returns kind=unsupported, no text): images, legacy .doc/.ppt/
    .xls binary formats, archives, and anything else. Caller should fall back to
    downloading the file (web_url / save_path).
    """
    import io

    lower = name.lower()
    ext = lower.rsplit(".", 1)[-1] if "." in lower else ""

    # 0. HTML first (its content_type is text/html, which would otherwise be
    #    swallowed by the plain-text catch-all below).
    if ext in ("html", "htm") or content_type == "text/html":
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        return ("html", soup.get_text(separator="\n", strip=True))

    # 1. Plain text / structured text.
    if ext in (
        "txt", "csv", "tsv", "md", "json", "xml", "log", "yaml", "yml", "ini", "conf"
    ) or content_type.startswith("text/"):
        text = raw.decode("utf-8", errors="replace")
        return ("csv" if ext in ("csv", "tsv") else "text", text)

    # 2. Excel + ODS spreadsheets (cells as tab-separated rows).
    if ext in ("xlsx", "xlsm"):
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"### Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if c is None else str(c) for c in row]
                if any(cells):
                    parts.append("\t".join(cells))
        wb.close()
        return ("xlsx", "\n".join(parts))

    if ext == "ods":
        from odf.opendocument import load as odf_load
        from odf.table import Table, TableRow, TableCell
        from odf.text import P as OdfP

        doc = odf_load(io.BytesIO(raw))
        parts = []
        for table in doc.getElementsByType(Table):
            parts.append(f"### Sheet: {table.getAttribute('name')}")
            for row in table.getElementsByType(TableRow):
                cells = []
                for cell in row.getElementsByType(TableCell):
                    cells.append(
                        " ".join(str(p) for p in cell.getElementsByType(OdfP))
                    )
                if any(cells):
                    parts.append("\t".join(cells))
        return ("xlsx", "\n".join(parts))

    # 3. PDF.
    if ext == "pdf" or content_type == "application/pdf":
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(raw))
        pages = [(p.extract_text() or "") for p in reader.pages]
        return ("pdf", "\n\n".join(f"--- page {i + 1} ---\n{t}" for i, t in enumerate(pages)))

    # 4. Word.
    if ext == "docx":
        import docx

        d = docx.Document(io.BytesIO(raw))
        return ("docx", "\n".join(p.text for p in d.paragraphs))

    if ext == "odt":
        from odf.opendocument import load as odf_load
        from odf.text import P as OdfP
        from odf import teletype

        doc = odf_load(io.BytesIO(raw))
        paras = [teletype.extractText(p) for p in doc.getElementsByType(OdfP)]
        return ("docx", "\n".join(paras))

    # 5. PowerPoint (slide text).
    if ext == "pptx":
        from pptx import Presentation

        prs = Presentation(io.BytesIO(raw))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"--- slide {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line:
                            parts.append(line)
        return ("pptx", "\n".join(parts))

    # 6. RTF.
    if ext == "rtf" or content_type == "application/rtf":
        from striprtf.striprtf import rtf_to_text

        return ("rtf", rtf_to_text(raw.decode("utf-8", errors="replace")))

    # 8. Email .eml (RFC822, stdlib).
    if ext == "eml" or content_type == "message/rfc822":
        import email
        from email import policy

        msg = email.message_from_bytes(raw, policy=policy.default)
        headers = "\n".join(
            f"{h}: {msg[h]}" for h in ("From", "To", "Cc", "Subject", "Date") if msg[h]
        )
        body_part = msg.get_body(preferencelist=("plain", "html"))
        body = body_part.get_content() if body_part else ""
        if body_part is not None and body_part.get_content_type() == "text/html":
            from bs4 import BeautifulSoup

            body = BeautifulSoup(body, "html.parser").get_text(separator="\n", strip=True)
        return ("eml", f"{headers}\n\n{body}")

    # 9. Outlook .msg.
    if ext == "msg":
        import extract_msg

        m = extract_msg.Message(io.BytesIO(raw))
        headers = "\n".join(
            f"{label}: {val}"
            for label, val in (
                ("From", m.sender),
                ("To", m.to),
                ("Cc", m.cc),
                ("Subject", m.subject),
                ("Date", m.date),
            )
            if val
        )
        body = m.body or ""
        m.close()
        return ("msg", f"{headers}\n\n{body}")

    return ("unsupported", "")


@mcp.tool
def get_event_attachment(
    event_id: str,
    attachment_id: str,
    account_id: str,
    save_path: str | None = None,
    onedrive_folder: str = "Attachments/Events",
) -> dict[str, Any]:
    """DOWNLOAD an event attachment to OneDrive/disk (NOT for reading contents).

    Use this only to save the file or get a shareable link. To READ what's inside
    a file (spreadsheet, PDF, etc), use read_event_attachment / read_attachment_text
    instead - they return the text directly.

    The attachment content is never returned inline as base64: a base64 blob in
    a tool result floods the LLM context and hangs the client (and the model
    cannot reconstruct binary formats anyway). Instead:

    - ``save_path`` set -> write the bytes to the server filesystem (local
      stdio only) and return ``saved_to``.
    - otherwise (HTTP/remote) -> upload the file to OneDrive under
      ``onedrive_folder`` and return ``onedrive_file_id`` + ``web_url``. The
      model gets a link it can hand to the user or open; the bytes stay out of
      the conversation.

    To let the model READ the contents (spreadsheet/pdf/etc), use
    ``read_event_attachment`` instead, which returns extracted text.

    Only ``fileAttachment`` types carry content; ``referenceAttachment`` (link
    attachments) have no bytes and raise.
    """
    name, content_type, raw = _fetch_event_attachment_bytes(
        event_id, attachment_id, account_id
    )
    response: dict[str, Any] = {
        "name": name,
        "content_type": content_type,
        "size": len(raw),
    }

    if save_path is not None:
        # Local stdio: persist to the server filesystem.
        path = pl.Path(save_path).expanduser().resolve()
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(raw)
        response["saved_to"] = str(path)
    else:
        # Remote/HTTP: stage to OneDrive, return a link (never the bytes).
        onedrive_path = f"{onedrive_folder.strip('/')}/{name}"
        uploaded = graph.upload_large_file(
            f"/me/drive/root:/{onedrive_path}:", raw, account_id
        )
        response["onedrive_path"] = onedrive_path
        response["onedrive_file_id"] = uploaded.get("id")
        response["web_url"] = uploaded.get("webUrl")
        response["note"] = (
            f"Saved to OneDrive at '{onedrive_path}'. Open web_url to view/"
            "download, or use read_event_attachment to read its contents as text."
        )

    return response


@mcp.tool
def read_event_attachment(
    event_id: str,
    attachment_id: str,
    account_id: str,
    max_chars: int = ATTACHMENT_TEXT_MAX_CHARS,
) -> dict[str, Any]:
    """USE THIS to read / open / view the CONTENTS of a calendar event attachment.

    This is the correct tool whenever the user asks to read, open, view, check,
    summarize or extract data from an event attachment (spreadsheet, PDF, ticket,
    document). It extracts the text server-side and returns it directly - do NOT
    use get_event_attachment for reading (that one only downloads the file).

    Extracts plain text on the server, so the model reads it without ever
    receiving base64. Supported:

    - xlsx/xlsm -> sheets and cells as tab-separated text
    - pdf       -> extracted text per page
    - docx      -> paragraph text
    - csv/tsv/txt/md/json/xml/log and any text/* -> decoded directly

    For unsupported binary types (e.g. images), no text is returned: use
    ``get_event_attachment`` to stage the file to OneDrive and open ``web_url``.

    Output text is truncated to ``max_chars`` to protect the context.
    """
    name, content_type, raw = _fetch_event_attachment_bytes(
        event_id, attachment_id, account_id
    )
    kind, text = _extract_text(name, content_type, raw)

    response: dict[str, Any] = {
        "name": name,
        "content_type": content_type,
        "size": len(raw),
        "kind": kind,
    }

    if kind == "unsupported":
        response["text"] = None
        response["note"] = (
            f"Cannot extract text from '{name}' ({content_type}). Use "
            "get_event_attachment to save it to OneDrive and open the link."
        )
        return response

    truncated = len(text) > max_chars
    response["text"] = text[:max_chars]
    response["truncated"] = truncated
    if truncated:
        response["note"] = (
            f"Text truncated to {max_chars} of {len(text)} chars. Raise "
            "max_chars or open the file in OneDrive (get_event_attachment) "
            "for the full content."
        )
    return response


def _build_text_response(
    name: str, content_type: str, raw: bytes, max_chars: int, download_hint: str
) -> dict[str, Any]:
    """Shared shaping for any 'read as text' tool."""
    kind, text = _extract_text(name, content_type, raw)
    response: dict[str, Any] = {
        "name": name,
        "content_type": content_type,
        "size": len(raw),
        "kind": kind,
    }
    if kind == "unsupported":
        response["text"] = None
        response["note"] = (
            f"Cannot extract text from '{name}' ({content_type}). {download_hint}"
        )
        return response
    truncated = len(text) > max_chars
    response["text"] = text[:max_chars]
    response["truncated"] = truncated
    if truncated:
        response["note"] = (
            f"Text truncated to {max_chars} of {len(text)} chars. Raise max_chars."
        )
    return response


def _download_onedrive_bytes(
    file_id: str, account_id: str
) -> tuple[str, str, bytes]:
    """Download a OneDrive item's bytes into memory via its pre-auth download URL.

    Returns (name, mime_type, raw). Uses @microsoft.graph.downloadUrl, which is
    short-lived and pre-authenticated (no Authorization header needed).
    """
    import httpx

    metadata = graph.request("GET", f"/me/drive/items/{file_id}", account_id)
    if not metadata:
        raise ValueError(f"File with ID {file_id} not found")
    download_url = metadata.get("@microsoft.graph.downloadUrl")
    if not download_url:
        raise ValueError("No download URL available for this file")
    resp = httpx.get(download_url, follow_redirects=True, timeout=60.0)
    resp.raise_for_status()
    name = metadata.get("name", "unknown")
    mime = (metadata.get("file", {}) or {}).get("mimeType", "application/octet-stream")
    return name, mime, resp.content


@mcp.tool
def read_attachment_text(
    account_id: str,
    event_id: str | None = None,
    email_id: str | None = None,
    attachment_id: str | None = None,
    onedrive_file_id: str | None = None,
    max_chars: int = ATTACHMENT_TEXT_MAX_CHARS,
) -> dict[str, Any]:
    """USE THIS to read / open / view the CONTENTS of any attachment or file.

    Preferred tool whenever the user wants to read, open, view, check, summarize
    or extract data from a file - event attachment, email attachment, or a
    OneDrive file. Returns the extracted text directly; never returns base64 and
    never needs the user to re-upload the file. Do NOT use the get_* tools for
    reading (those only download).

    One generic reader for the three sources - extracts text on the server (see
    _extract_text: text/csv/xlsx/ods/pdf/docx/odt/pptx/rtf/html/eml/msg).

    Provide exactly one source:
      - calendar event attachment:  event_id + attachment_id
      - email attachment:           email_id + attachment_id
      - OneDrive file:              onedrive_file_id

    For unsupported binary types (images, legacy .doc/.xls/.ppt, archives) the
    response has kind=unsupported and text=None; download the file instead
    (get_event_attachment / get_attachment / get_file).
    """
    if event_id and attachment_id:
        name, content_type, raw = _fetch_event_attachment_bytes(
            event_id, attachment_id, account_id
        )
        hint = "Use get_event_attachment to save it to OneDrive."
    elif email_id and attachment_id:
        result = graph.request(
            "GET", f"/me/messages/{email_id}/attachments/{attachment_id}", account_id
        )
        if not result:
            raise ValueError("Attachment not found")
        if "contentBytes" not in result:
            raise ValueError(
                "Attachment content not available - likely a reference (link) "
                "attachment with no downloadable bytes."
            )
        name = result.get("name", "unknown")
        content_type = result.get("contentType", "application/octet-stream")
        raw = base64.b64decode(result["contentBytes"])
        hint = "Use get_attachment to download it."
    elif onedrive_file_id:
        name, content_type, raw = _download_onedrive_bytes(onedrive_file_id, account_id)
        hint = "Use get_file to download it."
    else:
        raise ValueError(
            "Provide a source: (event_id + attachment_id), (email_id + "
            "attachment_id), or onedrive_file_id."
        )

    return _build_text_response(name, content_type, raw, max_chars, hint)


@mcp.tool
def search_files(
    query: str,
    account_id: str,
    limit: int = 50,
) -> list[dict[str, Any]]:
    """Search for files in OneDrive using the modern search API."""
    items = list(graph.search_query(query, ["driveItem"], account_id, limit))

    return [
        {
            "id": item["id"],
            "name": item["name"],
            "type": "folder" if "folder" in item else "file",
            "size": item.get("size", 0),
            "modified": item.get("lastModifiedDateTime"),
            "download_url": item.get("@microsoft.graph.downloadUrl"),
        }
        for item in items
    ]


@mcp.tool
def search_emails(
    query: str,
    account_id: str,
    limit: int = 50,
    folder: str | None = None,
) -> list[dict[str, Any]]:
    """Search emails by keyword.

    Uses the traditional ``$search`` endpoint on ``/me/messages`` (or a specific
    folder) rather than the modern ``/search/query`` API. The latter returns
    each hit's ``id`` in an OWA/EWS-style format that is NOT usable as a REST id
    by ``get_email`` / ``get_attachment`` (confirmed on the Brandart tenant),
    which breaks the search -> open -> download-attachment workflow. The
    ``$search`` endpoint returns the proper REST id via ``$select``.
    """
    if folder:
        folder_path = FOLDERS.get(folder.casefold(), folder)
        endpoint = f"/me/mailFolders/{folder_path}/messages"
    else:
        # Search the whole mailbox - same endpoint, no folder scope.
        endpoint = "/me/messages"

    params = {
        "$search": f'"{query}"',
        "$top": min(limit, 100),
        "$select": "id,subject,from,toRecipients,receivedDateTime,hasAttachments,body,conversationId,isRead",
    }

    return list(
        graph.request_paginated(endpoint, account_id, params=params, limit=limit)
    )


@mcp.tool
def search_events(
    query: str,
    account_id: str,
    days_ahead: int = 365,
    days_back: int = 365,
    limit: int = 50,
) -> list[dict[str, Any]]:
    """Search calendar events using the modern search API."""
    events = list(graph.search_query(query, ["event"], account_id, limit))

    # Filter by date range if needed
    if days_ahead != 365 or days_back != 365:
        now = dt.datetime.now(dt.timezone.utc)
        start = now - dt.timedelta(days=days_back)
        end = now + dt.timedelta(days=days_ahead)

        filtered_events = []
        for event in events:
            event_start = dt.datetime.fromisoformat(
                event.get("start", {}).get("dateTime", "").replace("Z", "+00:00")
            )
            event_end = dt.datetime.fromisoformat(
                event.get("end", {}).get("dateTime", "").replace("Z", "+00:00")
            )

            if event_start <= end and event_end >= start:
                filtered_events.append(event)

        return filtered_events

    return events


@mcp.tool
def search_contacts(
    query: str,
    account_id: str,
    limit: int = 50,
) -> list[dict[str, Any]]:
    """Search contacts. Uses traditional search since unified_search doesn't support contacts."""
    params = {
        "$search": f'"{query}"',
        "$top": min(limit, 100),
    }

    contacts = list(
        graph.request_paginated("/me/contacts", account_id, params=params, limit=limit)
    )

    return contacts


@mcp.tool
def search_people(
    query: str,
    account_id: str,
    limit: int = 25,
) -> list[dict[str, Any]]:
    """Search people in the organization directory (Global Address List) and frequent contacts.

    Uses the /me/people endpoint which searches across:
    - Azure AD / Entra directory (all company users)
    - Frequent contacts and communication patterns
    - Personal contacts

    Returns name, email, job title, department, and office location.
    """
    params = {
        "$search": f'"{query}"',
        "$top": min(limit, 100),
        "$select": "displayName,emailAddresses,jobTitle,department,officeLocation,companyName,userPrincipalName",
    }

    result = graph.request("GET", "/me/people", account_id, params=params)
    if not result or "value" not in result:
        return []

    people = []
    for person in result["value"]:
        emails = [
            e.get("address")
            for e in person.get("emailAddresses", [])
            if e.get("address")
        ]
        people.append(
            {
                "displayName": person.get("displayName"),
                "emails": emails,
                "jobTitle": person.get("jobTitle"),
                "department": person.get("department"),
                "officeLocation": person.get("officeLocation"),
                "companyName": person.get("companyName"),
            }
        )

    return people


@mcp.tool
def unified_search(
    query: str,
    account_id: str,
    entity_types: list[str] | None = None,
    limit: int = 50,
) -> dict[str, list[dict[str, Any]]]:
    """Search across multiple Microsoft 365 resources using the modern search API

    entity_types can include: 'message', 'event', 'drive', 'driveItem', 'list', 'listItem', 'site'
    If not specified, searches across all available types.
    """
    if not entity_types:
        entity_types = ["message", "event", "driveItem"]

    results = {entity_type: [] for entity_type in entity_types}

    items = list(graph.search_query(query, entity_types, account_id, limit))

    for item in items:
        resource_type = item.get("@odata.type", "").split(".")[-1]

        if resource_type == "message":
            results.setdefault("message", []).append(item)
        elif resource_type == "event":
            results.setdefault("event", []).append(item)
        elif resource_type in ["driveItem", "file", "folder"]:
            results.setdefault("driveItem", []).append(item)
        else:
            results.setdefault("other", []).append(item)

    return {k: v for k, v in results.items() if v}
